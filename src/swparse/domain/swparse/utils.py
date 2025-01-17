import io
import os
import re
import json
import math
import base64
import hashlib
import mimetypes
from uuid import uuid4
from itertools import islice
from operator import attrgetter
from datetime import UTC, datetime
from typing import Any, TYPE_CHECKING, List

 
import torch
import psutil
import structlog
import html_text
import mistletoe
import pandas as pd
import aioboto3
from botocore.exceptions import ClientError
from botocore.config import Config
from gliner import GLiNER
from lark import Lark, Token, Transformer
from litestar.exceptions import HTTPException, NotFoundException
from lxml import html
from markdown_it import MarkdownIt
from mdit_plain.renderer import RendererPlain
from nltk import download, tokenize
from openpyxl import load_workbook
from openpyxl.drawing.image import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import SlideShapes
from snakemd import Document as SnakeMdDocument
from snakemd.elements import MDList
from xls2xlsx import XLS2XLSX
from PIL import Image as PILImage

from swparse.config.app import settings
from swparse.db.models.content_type import ContentType
from swparse.domain.swparse.schemas import JobMetadata


if TYPE_CHECKING:
    from PIL.Image import Image

logger = structlog.get_logger()
BUCKET = settings.storage.BUCKET
JOB_FOLDER = settings.storage.JOB_FOLDER
MINIO_ROOT_USER = settings.storage.ROOT_USER
MINIO_ROOT_PASSWORD = settings.storage.ROOT_PASSWORD

SAQ_PROCESSES = settings.saq.PROCESSES

def convert_xls_to_xlsx_bytes(content: bytes) -> bytes:

    x2x = XLS2XLSX(io.BytesIO(content))
    workbook = x2x.to_xlsx()

    with io.BytesIO() as buffer:
        workbook.save(buffer)
        buffer.seek(0)
        return buffer.read()


def change_file_ext(file_name: str, extension: str) -> str:
    file_path = file_name.split(".")
    file_path[-1] = extension
    return ".".join(file_path)


async def extract_tables_from_html(html_file_path: str) -> list[str] | None:
    content = await read_file(html_file_path)
    content = content.decode()

    tree = html.fromstring(content)
    tables = tree.xpath("//table")
    if not tables:
        # No table found in the html file
        return None

    html_tables: list[str] = []
    for table in tables:
        html_table = html.tostring(table, method="html", pretty_print=True)

        if isinstance(html_table, bytes):
            html_table = html_table.decode()
        html_tables.append(html_table)

    return html_tables


def ungroup_shapes(shapes: SlideShapes) -> list[BaseShape]:
    res = []
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                res.extend(ungroup_shapes(shape.shapes))
            else:
                res.append(shape)
        except Exception as e:
            raise Exception(f"{e}")
    return res


def is_title(shape: BaseShape) -> bool:
    if shape.is_placeholder and shape.placeholder_format.type in [
        PP_PLACEHOLDER.TITLE,
        PP_PLACEHOLDER.SUBTITLE,
        PP_PLACEHOLDER.VERTICAL_TITLE,
        PP_PLACEHOLDER.CENTER_TITLE,
    ]:
        return True
    return False


def is_text_block(shape: BaseShape) -> bool:
    if shape.has_text_frame:
        if (
            (shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.BODY)
            or (shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.OBJECT)
            or shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
        ):
            return True
    return False


def is_list_block(shape: BaseShape) -> bool:
    levels = []
    for para in shape.text_frame.paragraphs:
        if para.level not in levels:
            levels.append(para.level)
        if para.level != 0 or len(levels) > 1:
            return True
    return False


def is_list_nested(shape: BaseShape) -> bool:
    levels = []
    for para in shape.text_frame.paragraphs:
        if para.level not in levels:
            levels.append(para.level)
        if len(levels) > 1:
            return True
    return False


def add_to_list(nested_list: MDList, level: int, text: str):
    if level == 0:
        nested_list._items.append(text)  # type: ignore
    else:
        if not nested_list._items or not isinstance(nested_list._items[-1], list):  # type: ignore
            nested_list._items.append(MDList([]))  # type: ignore
        add_to_list(nested_list._items[-1], level - 1, text)  # type: ignore


def process_shapes(shapes: list[BaseShape], file: SnakeMdDocument):
    try:
        for shape in shapes:
            if is_title(shape):
                file.add_heading(text=shape.text_frame.text, level=1)
            elif is_text_block(shape):
                if is_list_block(shape):
                    if is_list_nested(shape):
                        items = MDList([])
                        for paragraph in shape.text_frame.paragraphs:
                            add_to_list(items, paragraph.level, paragraph.text)
                        file.add_block(items)
                    else:
                        items = []
                        for paragraph in shape.text_frame.paragraphs:
                            add_to_list(items, paragraph.level, paragraph.text)
                        file.add_unordered_list(items=items)
                else:
                    for paragraph in shape.text_frame.paragraphs:
                        file.add_paragraph(text=paragraph.text)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pass
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                column_names = [cell.text for cell in shape.table.rows[0].cells]
                data = [[cell.text for cell in row.cells] for row in islice(shape.table.rows, 1, None)]
                file.add_table(
                    header=column_names,
                    data=data,
                )
            else:
                pass
            try:
                ph = shape.placeholder_format
                if ph.type == PP_PLACEHOLDER.OBJECT and hasattr(shape, "image") and getattr(shape, "image"):
                    pass
            except:
                pass
    except Exception as e:
        raise Exception


def convert_pptx_to_md(pptx_content: io.BytesIO, pptx_filename: str) -> str:
    try:
        prs = Presentation(pptx_content)
        md_file = SnakeMdDocument()
        for idx, slide in enumerate(prs.slides):
            shapes = []
            shapes = sorted(ungroup_shapes(slide.shapes), key=attrgetter("top", "left"))
            process_shapes(shapes, md_file)
        return str(md_file)
    except Exception as e:
        raise Exception





class MdAnalyser:
    def __init__(self, markdown_content: str):
        self.markdown_content = markdown_content
        self.components = []
        self.links = []
        self.lines = markdown_content.splitlines()

        # Compile patterns once for efficiency
        self.patterns = {
            "heading": re.compile(r"^(#{1,6})\s*(.+)$"),
            "table_row": re.compile(r"^\|.*\|$"),
            "paragraph": re.compile(r"^(?!\!\[.*\]\(.*\))(?!.*https?://)[^\|#\s].+$"),
            "table_separator": re.compile(r"^\|[-|]*\|$"),
            "links": re.compile(r"http?s?://[^\s]+")
        }

    def extract_components(self) -> tuple[List[dict[str,Any]], list[dict[str, str]]]:
        """Extracts components from the markdown content."""
        current_table = []

        for line in self.lines:
            if not line.strip():
                continue
            # Process different patterns
            self.extract_link(line)
            line = self.remove_links(line)
            if self.patterns["heading"].match(line):
                self._flush_table(current_table)
                self.add_heading(line)
            elif self.patterns["table_row"].match(line):
                if not self.patterns["table_separator"].match(line):
                    current_table.append(line)
            elif self.patterns["paragraph"].match(line):
                self._flush_table(current_table)
                self.add_paragraph(line)


        self._flush_table(current_table)

        return self.components, self.links


    def add_table(self, current_table: List[str]):
        """Processes and adds a table component."""
        rows = [
            [cell.strip() for cell in re.findall(r"\|([^|]+)", row)]
            for row in current_table
        ]
        self.components.append({
            "type": "table",
            "md": "\n".join(current_table),
            "rows": rows,
            "bBox": self._get_bbox(),
        })

    def _flush_table(self, current_table: List[str]):
        """Adds the current table to components if not empty."""
        if current_table:
            self.add_table(current_table)
            current_table.clear()


    def extract_link(self, line: str):
        """Extract links from the given line and adds them to self.links."""
        matches = self.patterns["links"].findall(line)
        for match in matches:
            self.links.append({
                "text": match,
                "url": match
            })

    def remove_links(self, line: str) -> str:
        """Remove links from a line while keeping the remaining text."""
        return self.patterns["links"].sub("", line).strip()

    def add_heading(self, line: str):
        """Processes and adds a heading component."""
        match = self.patterns["heading"].match(line)
        level = len(match.group(1))
        text = match.group(2).strip()
        self.components.append({
            "type": "heading",
            "level": level,
            "md": line,
            "value": text,
            "bBox": self._get_bbox()
        })

    def add_paragraph(self, line: str):
        """Processes and adds a text component."""
        self.components.append({
            "type": "text",
            "md": line,
            "value": self.md_to_text(line),
            "bBox": self._get_bbox()
        })

    @staticmethod
    def md_to_text(md: str) -> str:
        """Converts Markdown to plain text (basic implementation)."""
        return re.sub(r"(\*\*|__|`|~~)", "", md).strip()

    @staticmethod
    def _get_bbox() -> dict:
        """Returns a default bounding box structure."""
        return {"x": 0.0, "y": 0.0, "w": 0.0, "h": 0.0}



class TreeToJson(Transformer):
    def __init__(self, visit_tokens: bool = True):
        super().__init__(visit_tokens)
        self.tables = set()

    def start(self, items: list[Token]):
        return {"tables": items[0:-1], "output": items[-1].lower()}

    def instruction(self, items: list[Token]):
        headers = [label["name"] for label in items[1:-1]]
        if len(headers) != len(dict.fromkeys(headers)):
            raise ValueError("Duplicate labels")
        return {"mode": items[-1], "table_name": items[0], "labels": items[1:-1]}

    def mode(self, items: list[Token]):
        mode_map = {
            "by_sentence": "sent",
            "sent": "sent",
            "bysentence": "sent",
            "by_line": "ln",
            "ln": "ln",
            "byline": "ln",
        }
        if not len(items):
            return "sent"
        return mode_map[items[0].value.replace(" ", "_").lower()]

    def output(self, items: list[Token]):
        if not len(items):
            return "json"
        return items[0].value.lower()

    def table_ident(self, items: list[Token]):
        table_name = items[0].value.replace(" ", "_").lower()
        if table_name in self.tables:
            raise Exception(f"Duplicate table name: {table_name}")
        self.tables.add(table_name)
        return table_name

    def value(self, items: list[Token]):
        name = items[0].replace(" ", "_").lower()
        if len(items) == 1:
            return {"name": name, "type": "string"}
        return {"name": name, "type": items[1]}

    def field(self, items: list[Token]):
        return items[0].value

    def type(self, items: list[Token]):
        return items[0].value



def parse_table_query(extraction_query: str) -> list[dict]:
    lang = """start: instruction+ output
        instruction: table_ident value+ mode ";"
        table_ident: FIELD_NAME "="
        value: field type?
        field: FIELD_NAME ","?
        mode: ("-"MODE)?
        output: (("as "|"AS ")OUTPUT)?
        OUTPUT: /csv|md|json|html/i
        MODE: /by( |_)?sentence|by( |_)?line|ln|sent|sentence|line/i
        FIELD_NAME: /[a-zA-Z0-9_]+((_| )[a-zA-Z0-9_]+)*/
        type: ":" DATATYPES ","?
        DATATYPES: DATATYPE"[]"?
        DATATYPE: "str" | "string" | "text" | "int" | "integer" | "number" | "float" | "date" | "bool" | "boolean"
        %import common.WS
        %ignore WS
        """
    lang_parser = Lark(lang)
    tree = lang_parser.parse(extraction_query)
    return TreeToJson().transform(tree)


def extract_tables_gliner(table_queries: list[dict], markdownText: str, output: str | None = None) -> Any:
    download("punkt_tab")
    model = GLiNER.from_pretrained("gliner-community/gliner_large-v2.5")

    text = html_text.extract_text(mistletoe.markdown(markdownText))
    sent_tokens = tokenize.sent_tokenize(text)
    ln_tokens = tokenize.line_tokenize(text)

    results = []

    for query in table_queries:
        list_labels = []
        int_labels = []
        float_labels = []
        labels = []
        for label in query["labels"]:
            if label["type"][-2:] == "[]":
                list_labels.append(label["name"])
                label["type"] = label["type"][:-2]
            if label["type"] in ["int", "integer", "number"]:
                int_labels.append(label["name"])
            elif label["type"] in ["float"]:
                float_labels.append(label["name"])
            labels.append(label["name"])

        table = {
            "table_name": query["table_name"],
            "headers": labels,
            "rows": [],
        }

        tokens = sent_tokens if query["mode"] == "sent" else ln_tokens

        for token in tokens:
            entities = model.predict_entities(token, labels, threshold=0.5)
            row = {}
            for list_label in list_labels:
                row[list_label] = []

            extracted_labels = list(dict.fromkeys([entity["label"] for entity in entities]))

            if len(extracted_labels) < (math.floor(len(labels) * 0.8) if len(labels) > 1 else 1):
                continue

            for entity in entities:
                if entity["label"] in int_labels or entity["label"] in float_labels:
                    match = re.search(r"\b\d{1,3}(?:,\d{3})*(?:\.\d+)?\b", entity["text"])
                    if entity["label"] in int_labels:
                        entity["text"] = int(match.group()) if match else None
                    else:
                        entity["text"] = float(match.group().replace(",", "")) if match else None

                if entity["text"]:
                    if entity["label"] in list_labels:
                        row[entity["label"]].append(entity["text"])
                    else:
                        row[entity["label"]] = entity["text"]

            table["rows"].append(row)

        results.append(table)

    data_frames = []
    table_names = []
    max_header = 0
    for table in results:
        rows = []
        for row in table["rows"]:
            processed_row = {col: ", ".join(value) if isinstance(value, list) else value for col, value in row.items()}
            rows.append(processed_row)

        if max_header < len(table["headers"]):
            max_header = len(table["headers"])
        df = pd.DataFrame(rows, columns=table["headers"])

        data_frames.append(df.fillna(""))
        table_names.append(table["table_name"])

    json_result = json.dumps(results)
    csv_result = ""
    md_result = ""

    for i, df in enumerate(data_frames):
        # MARKDOWN
        md_result += f"# {table_names[i]}\n"
        md_result += df.to_markdown(index=False) + "\n"

        # CSV
        csv_buffer = io.StringIO()
        df.to_csv(csv_buffer, index=False)
        csv_string = csv_buffer.getvalue()

        if i != 0:
            csv_result += f"\n============={',' * (max_header - 1)}\n\n"

        csv_result += f"{table_names[i]}{',' * (max_header - 1)}\n"
        csv_result += csv_string

    # HTML
    html_result = mistletoe.markdown(md_result)

    if output == "json":
        return json_result
    elif output == "csv":
        return csv_result
    elif output == "md":
        return md_result
    elif output == "html":
        return html_result

    if output is None:
        return {
            "csv": csv_result,
            "md": md_result,
            "html": html_result,
            "json": json_result,
        }


def get_hashed_file_name(filename: str, hashed_input: dict[str, Any]) -> str:
    """
    Generate a hashed file name based on file content or additional input data.

    Args:
        filename (str): Original file name.
        input (Union[bytes, dict]): File content as bytes or a dictionary with additional instruction fields.

    Returns:
        str: Hashed file name.
    """
    file_ext = filename.split(".")[-1]

    key_len = len(list(hashed_input.keys()))

    if key_len == 1:
        input_bytes = hashed_input["content"]

    else:
        hashed_input["content"] = base64.b64encode(hashed_input["content"]).decode("utf-8")

        if hashed_input.get("sheet_index"):
            hashed_input["sheet_index"] = sorted(hashed_input["sheet_index"], key=lambda x: str(x))

        if hashed_input.get("force_ocr"):
            hashed_input["force_ocr"] = True

        input_bytes = json.dumps(hashed_input, sort_keys=True).encode("utf-8")

    check_sum = hashlib.md5(input_bytes).hexdigest()
    return f"{check_sum}.{file_ext}"



async def get_file_content(file_path: str) -> str:
    content = await read_file(file_path)
    if isinstance(content, bytes):
        content = content.decode()
    return content


async def handle_result_type(
    result_type: str, results: dict[str, str], jm: JobMetadata, job_key: str
) -> dict:
    if results.get("result_type"):
        result_type_check = results["result_type"]
    else:
        result_type_check = result_type
    try:
        result = {"job_metadata": jm.__dict__}
        if result_type_check == ContentType.MARKDOWN.value:
            markdown = await get_file_content(results["markdown"])
            result[result_type] = markdown

        elif result_type_check == ContentType.HTML.value:
            html = await get_file_content(results["html"])
            result[result_type] = html

        elif result_type_check == ContentType.TEXT.value:
            text =  await get_file_content(results["text"])
            result[result_type] = text

        elif result_type_check == ContentType.TABLE.value:
            html = await extract_tables_from_html(results["html"])
            result_html = "<br><br>"
            if html:
                result_html = result_html.join(html)
            result[result_type] = result_html

        elif result_type_check == ContentType.JSON.value:
            json_str =  await get_file_content(results["json"])
            json_pages = json.loads(json_str)
            result = {
                'pages':json_pages,
                'job_metadata': jm.__dict__,
            }

        elif result_type_check == ContentType.MARKDOWN_TABLE.value:
            table_file_path = results.get("table")
            if table_file_path:
                 
                result_html = await get_file_content(results["table"])
                result_html = result_html.decode()
            else:
                html = await extract_tables_from_html(results["html"])
                result_html = "<br><br>".join(html)

            dfs = pd.read_html(result_html)
            markdown_tbls = ""
            for i, df in enumerate(dfs):
                markdown_tbls += f"## Table {i + 1}\n\n"
                markdown_tbls += df.to_markdown()
                markdown_tbls += "\n\n"

            result[result_type] = markdown_tbls

        elif result_type_check == ContentType.IMAGES.value:
           
            json_image_meta = await get_file_content(results["images"])

            result[result_type] = json.loads(json_image_meta)
        else:
            tables_content = await get_file_content(results[result_type_check])
            result[result_type] = tables_content

        if result:
            return result

        raise HTTPException(f"Format {result_type_check} is currently unsupported for {job_key}")

    except Exception as e:
        raise HTTPException(f"Format {result_type_check} is currently unsupported for {job_key}")


def md_to_text(md: str) -> str:
    parser = MarkdownIt(renderer_cls=RendererPlain)
    return parser.render(md)



def extract_md_components(markdown_content: str)->tuple[list[dict[str, Any]], list[dict[str, str]]]:
    analyser = MdAnalyser(markdown_content)
    return analyser.extract_components()


async def extract_excel_images(excel_content: io.BytesIO, sheet_name: str|int) -> dict[str, str]:
    """
    Extract images from an Excel file and store them in S3.
    return a dictionary mapping image names to S3 paths.
    """
    pxl_doc = load_workbook(filename= excel_content)
    sheet_images = {}

    # openpyxl only work with sheet name
    if isinstance(sheet_name, int):
        # retrieve sheet name based on the sheet index
        sheet_name = pxl_doc.sheetnames[sheet_name]

    sheet = pxl_doc[sheet_name]

    for i, image in enumerate(sheet._images):
        if isinstance(image, Image):
            img_data = image.ref
            pil_image = PILImage.open(img_data)
            img_name = f"sheet_{sheet_name.lower()}_image_{i}.png"

            saved_img_path = await save_image(img_name, pil_image)

            sheet_images[img_name] = saved_img_path

    return sheet_images


def format_timestamp(timestamp:float) ->str:
    value = datetime.fromtimestamp(timestamp)
    return value.strftime('%S:') + f"{int(value.strftime('%f')) // 1000}"


def get_memory_usage():
    process = psutil.Process(os.getpid())
    return process.memory_info()

def get_vram_usage():
    if torch.cuda.is_available():
        allocated = torch.cuda.memory_allocated() / 1024**2
        cached = torch.cuda.memory_reserved() / 1024**2
        return allocated, cached



async def get_file_name(s3_url: str) -> str:
    return os.path.basename(s3_url).split("/")[-1]



def parse_minio_url(s3_url: str):
    """Parse MinIO-style S3 URL into bucket name and key i.e compatible with aws s3."""
    parts = s3_url.split("/", 1)
    if len(parts) != 2:
        raise ValueError("Invalid MinIO URL format. Expected 'BUCKET/KEY'.")
    return parts[0], parts[1]

 

# def parse_s3_url(s3_url: str):
#     """Parse Minio S3 URL into bucket name and key."""
#     if not s3_url.startswith("s3://"):
#         raise ValueError("Invalid S3 URL")
#     parts = s3_url[5:].split("/", 1)
#     return parts[0], parts[1]



async def read_file(s3_url: str) -> bytes | str:
    """Asynchronously read file content from MinIO."""

    bucket_name, key = parse_minio_url(s3_url)

    session = aioboto3.Session()

    async with session.client(
        "s3",
        endpoint_url=settings.storage.ENDPOINT_URL,
        aws_access_key_id=MINIO_ROOT_USER,
        aws_secret_access_key=MINIO_ROOT_PASSWORD,
        config=Config(signature_version="s3v4"),
    ) as s3_client:

        response = await s3_client.get_object(Bucket=bucket_name, Key=key)
        content = await response["Body"].read()
        return content


async def save_file(filename: str, content: bytes, randomize:bool = True) -> str:
    """Asynchronously save file content to MinIO."""
   
    if randomize:
        new_uuid = uuid4()
        key = f"{new_uuid}_{filename}"
    else:
        key = filename
        
    s3_url = f"{BUCKET}/{key}"
    mime_type, _ = mimetypes.guess_type(filename)

    session = aioboto3.Session()

    async with session.client(
        "s3",
        endpoint_url=settings.storage.ENDPOINT_URL,
        aws_access_key_id=MINIO_ROOT_USER,
        aws_secret_access_key=MINIO_ROOT_PASSWORD,
        config=Config(signature_version="s3v4"),
    ) as s3_client:

        await s3_client.put_object(Bucket=BUCKET, Key=key, Body=content, ContentType=mime_type)
        return s3_url


def save_image_sync(s3_client: Any, image_name: str, image: "PILImage") -> str:
    """Synchronous wrapper around async_save_image."""
 
    image_name = image_name.lower()
    buffered = io.BytesIO()
    image.save(buffered, format=image_name.split(".")[-1])
    img_b = buffered.getvalue()
  
    new_uuid = uuid4()
 
    key = f"{new_uuid}_{image_name}"
    s3_url = f"{BUCKET}/{key}"
    mime_type, _ = mimetypes.guess_type(image_name)

    try:
        s3_client.put_object(Bucket=BUCKET, Key=key, Body=img_b, ContentType=mime_type)
        return s3_url
    except ClientError as e:
        logger.error(f"Error uploading image: {e}")
        raise Exception(f"Failed to save image: {e}")
    
async def save_image(image_name: str, image: "PILImage") -> str:
    # Lowercase the image name for consistency.
    image_name = image_name.lower()

    buffered = io.BytesIO()
    image.save(buffered, format=image_name.split(".")[-1])
    img_b = buffered.getvalue()

    return await save_file(image_name, img_b)


# Meta data utils

async def save_metadata(s3_url: str, metadata: dict) -> None:
    """Asynchronously save metadata to a separate S3 object."""
   
    bucket_name, key = s3_url.split("/", 1)
    session = aioboto3.Session()
    try:
        async with session.client(
            "s3",
            endpoint_url=settings.storage.ENDPOINT_URL,
            aws_access_key_id=MINIO_ROOT_USER,
            aws_secret_access_key=MINIO_ROOT_PASSWORD,
            config=Config(signature_version="s3v4"),
        ) as s3_client:
            response = await s3_client.get_object(Bucket=bucket_name, Key=key)
            content = await response["Body"].read()
            
            if content is None:
                raise NotFoundException(detail=f"There is no {s3_url}")
            
            await s3_client.put_object(
                Bucket=bucket_name,
                Key=key,
                Body=content,
                ContentType="application/json",
                Metadata=metadata
            )
    except Exception as e:
        logger.info(e)
        logger.info("error when working with")
        # logger.info(metadata)
        logger.info(type(metadata))


 

async def get_job_metadata(job_id: str) -> dict[str, dict[str, str]]:
    """Asynchronously fetch job metadata from an S3 object."""
    key = f"{JOB_FOLDER}/{job_id}.json"

    session = aioboto3.Session()
    async with session.client(
        "s3",
        endpoint_url=settings.storage.ENDPOINT_URL,
        aws_access_key_id=MINIO_ROOT_USER,
        aws_secret_access_key=MINIO_ROOT_PASSWORD,
        config=Config(signature_version="s3v4"),
    ) as s3_client:
        try:
            response = await s3_client.get_object(Bucket=BUCKET, Key=key)
            content = await response["Body"].read()
            return json.loads(content)
        except s3_client.exceptions.NoSuchKey:
            # metatadata not found
            return {}

 
async def save_job_metadata(job_id: str, metadata: dict[str, str]) -> dict[str, dict[str, str]]:
    """Save job metadata asynchronously to an S3-compatible storage."""
    bucket = BUCKET
    job_file_key = f"{JOB_FOLDER}/{job_id}.json"
    job_folder_key = f"{JOB_FOLDER}/"

    existing_job_metadata: dict[str, dict] = {"metadata": {}}

    session = aioboto3.Session()

    async with session.client(
        "s3",
        endpoint_url=settings.storage.ENDPOINT_URL,
        aws_access_key_id=MINIO_ROOT_USER,
        aws_secret_access_key=MINIO_ROOT_PASSWORD,
        config=Config(signature_version="s3v4"),
    ) as s3_client:
        # Check if folder exists 
        try:
            response = await s3_client.list_objects_v2(Bucket=bucket, Prefix=job_folder_key, MaxKeys=1)
            folder_exists = "Contents" in response
        except Exception as err:
            raise HTTPException(status_code=500, detail=f"Error checking folder existence: {err}")

        if not folder_exists:
           
            try:
                await s3_client.put_object(Bucket=bucket, Key=f"{job_folder_key}", Body="")
            except Exception as err:
                raise HTTPException(status_code=500, detail=f"Error creating folder: {err}")

        try:
            await s3_client.head_object(Bucket=bucket, Key=job_file_key)
            # File exists; fetch existing metadata
            response = await s3_client.get_object(Bucket=bucket, Key=job_file_key)
            content = await response["Body"].read()
            existing_job_metadata = json.loads(content)
            logger.info("existing_job_metadata")
            logger.info(existing_job_metadata)
        except Exception as e:
            pass
            
        # File does not exist; create new metadata
        if not existing_job_metadata:
            existing_job_metadata["metadata"] = {"created_at": datetime.now(UTC).isoformat()}
            
        existing_job_metadata["metadata"].update(metadata)

        # Save updated metadata
        try:
            await s3_client.put_object(
                Bucket=bucket,
                Key=job_file_key,
                Body=json.dumps(existing_job_metadata, indent=4),
            )
            return existing_job_metadata
        except Exception as err:
            raise HTTPException(status_code=500, detail=f"Error saving job metadata: {err}")


async def is_file_exist(s3_url: str) -> bool:
    """Check if a file exists in S3/MinIO."""
    bucket_name, key = s3_url.split("/", 1)

    session = aioboto3.Session()

    async with session.client(
        "s3",
        endpoint_url=settings.storage.ENDPOINT_URL,
        aws_access_key_id=MINIO_ROOT_USER,
        aws_secret_access_key=MINIO_ROOT_PASSWORD,
        config=Config(signature_version="s3v4"),
    ) as s3_client:
        try:
         
            await s3_client.head_object(Bucket=bucket_name, Key=key)
            return True
        except Exception as e:
            ## File does not exist
            return False
         


async def get_metadata(s3_url: str) -> dict[str, str]:
    """Asynchronously retrieve metadata from an S3 object."""
    
    bucket_name, key = s3_url.split("/", 1)
    session = aioboto3.Session()
    try:
        async with session.client(
            "s3",
            endpoint_url=settings.storage.ENDPOINT_URL,
            aws_access_key_id=MINIO_ROOT_USER,
            aws_secret_access_key=MINIO_ROOT_PASSWORD,
            config=Config(signature_version="s3v4"),
        ) as s3_client:
  
            response = await s3_client.head_object(Bucket=bucket_name, Key=key)
            
  
            metadata = response.get("Metadata", {})
            if isinstance(metadata, str):
                return json.loads(metadata)
            return metadata
    except Exception as e:
        logger.error(e)
        raise NotFoundException(detail=f"Object not found: {s3_url}")
  